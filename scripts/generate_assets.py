from __future__ import annotations

import argparse
from pathlib import Path
import yaml

ROOT = Path(__file__).resolve().parents[1]


def main() -> None:
    parser = argparse.ArgumentParser(description="Genera PDF/PPTX ligeros desde el currículo")
    parser.add_argument("--lesson", help="ID de clase; omitir para todas")
    args = parser.parse_args()
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit('instala extras: pip install -e ".[assets]"') from exc

    curriculum = yaml.safe_load((ROOT / "curriculum.yaml").read_text(encoding="utf-8"))
    lessons = [item for part in curriculum["parts"] for item in part["lessons"]]
    if args.lesson:
        lessons = [item for item in lessons if item["id"] == args.lesson.zfill(3)]
    out = ROOT / "docs" / "generated"
    out.mkdir(parents=True, exist_ok=True)
    for item in lessons:
        pdf = canvas.Canvas(str(out / f"{item['id']}.pdf"), pagesize=A4)
        pdf.setFont("Helvetica-Bold", 16)
        pdf.drawString(48, 790, f"{item['id']} - {item['title']}"[:90])
        pdf.setFont("Helvetica", 10)
        pdf.drawString(48, 760, item["summary"][:100])
        pdf.save()

        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[0])
        slide.shapes.title.text = f"{item['id']} - {item['title']}"
        slide.placeholders[1].text = item["summary"]
        deck.save(out / f"{item['id']}.pptx")
    print(f"generados {len(lessons)} PDF y {len(lessons)} PPTX en {out}")


if __name__ == "__main__":
    main()
